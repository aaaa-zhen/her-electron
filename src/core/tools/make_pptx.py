#!/usr/bin/env python3
"""Generate a polished PPTX from JSON input.

Usage: echo '<json>' | python3 make_pptx.py <output_path>

JSON schema:
{
  "title": "Presentation Title",
  "subtitle": "Optional subtitle",
  "author": "Author name",
  "theme": "dark" | "light" | "blue" | "green",
  "slides": [
    {
      "layout": "title" | "content" | "two_column" | "image" | "quote" | "blank",
      "title": "Slide Title",
      "body": "Body text or bullet points separated by \\n",
      "left": "Left column text (for two_column)",
      "right": "Right column text (for two_column)",
      "quote": "Quote text (for quote layout)",
      "author": "Quote author",
      "notes": "Speaker notes"
    }
  ]
}
"""

import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

THEMES = {
    "dark": {
        "bg": RGBColor(0x1a, 0x1a, 0x2e),
        "title_color": RGBColor(0xff, 0xff, 0xff),
        "body_color": RGBColor(0xcc, 0xcc, 0xcc),
        "accent": RGBColor(0x6e, 0xe7, 0xb7),
        "subtitle_color": RGBColor(0x99, 0x99, 0xaa),
        "shape_fill": RGBColor(0x25, 0x25, 0x3d),
    },
    "light": {
        "bg": RGBColor(0xf8, 0xf8, 0xfc),
        "title_color": RGBColor(0x1a, 0x1a, 0x2e),
        "body_color": RGBColor(0x44, 0x44, 0x55),
        "accent": RGBColor(0x10, 0xb9, 0x81),
        "subtitle_color": RGBColor(0x77, 0x77, 0x88),
        "shape_fill": RGBColor(0xee, 0xee, 0xf4),
    },
    "blue": {
        "bg": RGBColor(0x0f, 0x17, 0x2a),
        "title_color": RGBColor(0xff, 0xff, 0xff),
        "body_color": RGBColor(0xbb, 0xcc, 0xdd),
        "accent": RGBColor(0x60, 0xa5, 0xfa),
        "subtitle_color": RGBColor(0x88, 0x99, 0xbb),
        "shape_fill": RGBColor(0x1e, 0x29, 0x3b),
    },
    "green": {
        "bg": RGBColor(0x0a, 0x1a, 0x15),
        "title_color": RGBColor(0xff, 0xff, 0xff),
        "body_color": RGBColor(0xbb, 0xdd, 0xcc),
        "accent": RGBColor(0x34, 0xd3, 0x99),
        "subtitle_color": RGBColor(0x88, 0xbb, 0x99),
        "shape_fill": RGBColor(0x15, 0x2e, 0x22),
    },
}


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=RGBColor(0xff, 0xff, 0xff), bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, text, font_size=16,
                    color=RGBColor(0xcc, 0xcc, 0xcc), accent=RGBColor(0x6e, 0xe7, 0xb7)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        is_bullet = line.startswith("- ") or line.startswith("* ") or line.startswith("• ")
        if is_bullet:
            line = line.lstrip("-*• ").strip()
            p.text = f"  •  {line}"
        else:
            p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.line_spacing = Pt(font_size * 1.6)
    return txBox


def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = 0.0
    return shape


def make_title_slide(prs, slide_data, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, theme["bg"])

    # Accent bar at top
    add_accent_bar(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), theme["accent"])

    # Title
    title = slide_data.get("title", "")
    add_text_box(slide, Inches(1), Inches(2.2), Inches(8), Inches(1.5),
                 title, font_size=40, color=theme["title_color"], bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Subtitle
    subtitle = slide_data.get("body", "") or slide_data.get("subtitle", "")
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(3.8), Inches(7), Inches(0.8),
                     subtitle, font_size=20, color=theme["subtitle_color"],
                     alignment=PP_ALIGN.CENTER)

    # Bottom accent line
    add_accent_bar(slide, Inches(4), Inches(3.55), Inches(2), Inches(0.04), theme["accent"])

    return slide


def make_content_slide(prs, slide_data, theme, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    # Left accent bar
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), theme["accent"])

    # Slide number
    add_text_box(slide, Inches(0.3), Inches(0.3), Inches(0.6), Inches(0.4),
                 f"{slide_num:02d}", font_size=11, color=theme["accent"], bold=True)

    # Title
    title = slide_data.get("title", "")
    if title:
        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
                     title, font_size=28, color=theme["title_color"], bold=True)
        # Underline
        add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(1.2), Inches(0.03), theme["accent"])

    # Body
    body = slide_data.get("body", "")
    if body:
        add_bullet_text(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5),
                        body, font_size=18, color=theme["body_color"], accent=theme["accent"])

    return slide


def make_two_column_slide(prs, slide_data, theme, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.06), Inches(7.5), theme["accent"])
    add_text_box(slide, Inches(0.3), Inches(0.3), Inches(0.6), Inches(0.4),
                 f"{slide_num:02d}", font_size=11, color=theme["accent"], bold=True)

    title = slide_data.get("title", "")
    if title:
        add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.8),
                     title, font_size=28, color=theme["title_color"], bold=True)
        add_accent_bar(slide, Inches(0.8), Inches(1.15), Inches(1.2), Inches(0.03), theme["accent"])

    # Left column bg
    left_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.6), Inches(1.6), Inches(4.1), Inches(5))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = theme["shape_fill"]
    left_shape.line.fill.background()

    # Right column bg
    right_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(5.1), Inches(1.6), Inches(4.1), Inches(5))
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = theme["shape_fill"]
    right_shape.line.fill.background()

    left_text = slide_data.get("left", "")
    right_text = slide_data.get("right", "")

    if left_text:
        add_bullet_text(slide, Inches(0.9), Inches(1.9), Inches(3.5), Inches(4.4),
                        left_text, font_size=16, color=theme["body_color"])
    if right_text:
        add_bullet_text(slide, Inches(5.4), Inches(1.9), Inches(3.5), Inches(4.4),
                        right_text, font_size=16, color=theme["body_color"])

    return slide


def make_quote_slide(prs, slide_data, theme, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    # Large quote mark
    add_text_box(slide, Inches(1), Inches(1.5), Inches(1), Inches(1),
                 "\u201C", font_size=72, color=theme["accent"], bold=True)

    quote = slide_data.get("quote", "") or slide_data.get("body", "")
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(7), Inches(3),
                 quote, font_size=24, color=theme["title_color"],
                 alignment=PP_ALIGN.LEFT)

    author = slide_data.get("author", "")
    if author:
        add_text_box(slide, Inches(1.5), Inches(5.2), Inches(7), Inches(0.5),
                     f"\u2014 {author}", font_size=16, color=theme["accent"])

    return slide


def generate(data, output_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    theme_name = data.get("theme", "dark")
    theme = THEMES.get(theme_name, THEMES["dark"])

    # Cover slide
    cover = {
        "title": data.get("title", "Untitled"),
        "body": data.get("subtitle", ""),
    }
    make_title_slide(prs, cover, theme)

    slide_num = 1
    for s in data.get("slides", []):
        layout = s.get("layout", "content")
        if layout == "title":
            make_title_slide(prs, s, theme)
        elif layout == "two_column":
            make_two_column_slide(prs, s, theme, slide_num)
        elif layout == "quote":
            make_quote_slide(prs, s, theme, slide_num)
        else:
            make_content_slide(prs, s, theme, slide_num)
        slide_num += 1

        # Speaker notes
        notes = s.get("notes", "")
        if notes:
            slide = prs.slides[-1]
            if slide.has_notes_slide:
                slide.notes_slide.notes_text_frame.text = notes

    prs.save(output_path)
    print(f"OK:{output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: echo '<json>' | python3 make_pptx.py <output>", file=sys.stderr)
        sys.exit(1)
    raw = sys.stdin.read()
    data = json.loads(raw)
    generate(data, sys.argv[1])
